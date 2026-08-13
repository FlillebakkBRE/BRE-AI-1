#!/usr/bin/env python3
"""Legger til en «Flaskehals-status»-slide i styremøte-decket, rett etter handlingsplanen.
Idempotent: fjerner ev. tidligere flaskehals-slide (markert via title-tekst) før ny bygges."""
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "BRE_Styremote_Finansiering_2026_18.pptx"
prs = Presentation(OUT)

BLA=RGBColor(0x00,0x56,0x89); MID=RGBColor(0x00,0x92,0xD2); LYS=RGBColor(0x59,0xC2,0xEA)
GUL=RGBColor(0xFA,0xE1,0x00); HVIT=RGBColor(0xFF,0xFF,0xFF); MORK=RGBColor(0x22,0x2B,0x33)
GRA=RGBColor(0xEE,0xF2,0xF5); GRAA=RGBColor(0x88,0x93,0x9b); LB=RGBColor(0xDD,0xEE,0xF6)
GRONN=RGBColor(0x2E,0xA0,0x4B); ROD=RGBColor(0xD8,0x3A,0x3A); GUL_D=RGBColor(0xE6,0xB0,0x00)

MARK = "FLASKEHALS-STATUS-SLIDE-BRE"

# --- idempotent: fjern tidligere flaskehals-slide om den finnes ---
def slide_has_mark(s):
    for sh in s.shapes:
        if sh.has_text_frame and MARK in sh.text_frame.text:
            return True
    return False
lst = prs.slides._sldIdLst
for sid in list(lst):
    rId = sid.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    part = prs.part.related_part(rId)
    s = next((sl for sl in prs.slides if sl.part is part), None)
    if s and slide_has_mark(s):
        lst.remove(sid); prs.part.drop_rel(rId)

def blank_slide():
    lay=min(prs.slide_layouts,key=lambda L:len(L.placeholders))
    s=prs.slides.add_slide(lay)
    for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
    return s

s = blank_slide()

def box(l,t,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh

def tx(l,t,w,h,text,size=12,color=MORK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,name="Calibri"):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name=name
    return tb

# tittelbånd
box(0,0,13.333,0.9,fill=BLA)
tx(0.5,0.12,12.3,0.4,"Nøkkelutfordringer / flaskehalser — status nå",21,HVIT,True)
tx(0.5,0.54,12.3,0.3,"Sensacon-rapportens seks flaskehalser vurdert per august 2026 (etter Roger + Even, Admento på stab, sterk H1)",11,LYS,False)
tx(11.9,0.0,1.4,0.02,MARK,1,BLA)  # skjult markør (bitteliten, i tittelfargen)

items=[
 (GUL_D,"Rekrutteringstilgang / kapasitet",
  "Bedret — Roger (salg) og Even (utvikling) inn, F. Johansen frigjort til leveranse.\nLeveransekapasitet er fortsatt hovedbremsen; software #2 + fagkapasitet gjenstår.","Delvis"),
 (GUL_D,"Salg + partnerskap / verdiskaping",
  "CSO på plass og HubSpot-pipeline gir struktur i salgsarbeidet.\nPartnerprogram og tydeligere synliggjøring av verdiskaping gjenstår.","Delvis"),
 (GUL_D,"Segmentprofil / prioritering",
  "Prioritering satt: logistikk, datasenter, energi.\nIkke operasjonalisert i daglig salg ennå — kommer når CSO jobber segmentspesifikt.","Delvis"),
 (ROD,"Modulbasert produktportefølje",
  "Ingen modularisering gjennomført ennå — brems på skalerbarhet.\nDekket av handlingsplanens produktgjennomgang («80/20 standard/opsjon»).","Åpen"),
 (GRONN,"Administrative oppgaver på nøkkelpersonell",
  "I stor grad løst — Admento tar personal + økonomi eksternt og avlaster DL.\nEgen controller ligger som neste forsterkning.","Løst"),
 (ROD,"Prismodell",
  "Uklar og lite strukturert — reell svakhet i salg og marginstyring.\nForenkling planlagt som tiltak i Q4 (internt + eksternt formål).","Åpen"),
]

# 2 kolonner x 3 rader
colx=[0.55,6.95]; cardw=5.83; cardh=1.62
rowy=[1.15,2.93,4.71]
for i,(dot,title,body,tag) in enumerate(items):
    x=colx[i%2]; y=rowy[i//2]
    box(x,y,cardw,cardh,fill=GRA,line=LYS,lw=1.0)
    # statusprikk
    box(x+0.22,y+0.24,0.42,0.42,fill=dot,shape=MSO_SHAPE.OVAL)
    tx(x+0.8,y+0.16,cardw-1.6,0.4,title,13,BLA,True)
    # status-tag (høyre)
    box(x+cardw-1.15,y+0.2,0.95,0.4,fill=dot)
    tx(x+cardw-1.15,y+0.2,0.95,0.4,tag,10,HVIT,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    tx(x+0.8,y+0.62,cardw-1.0,cardh-0.72,body,10.5,MORK,False)

# oppsummeringsbånd nederst
by=6.55
box(0.55,by,12.23,0.62,fill=BLA)
tx(0.75,by+0.08,12.0,0.46,"Kortversjon:  1 løst (admin)  ·  3 på vei (kapasitet, salg, segment)  ·  2 gjenstår (produkt + prismodell) — alle med tiltak i handlingsplanen.",12.5,HVIT,True,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)

# --- plasser rett etter handlingsplanen (slide 15 => idx 14 => ny blir idx 15) ---
ids=list(lst)
new_node=ids[-1]
lst.remove(new_node)
lst.insert(15, new_node)

prs.save(OUT)
print("Ferdig. Antall slides:", len(prs.slides._sldIdLst))
# verifiser plassering
prs2=Presentation(OUT)
for i,sl in enumerate(prs2.slides):
    if slide_has_mark(sl):
        print(f"Flaskehals-slide ligger på indeks {i} (slide {i+1}).")
