#!/usr/bin/env python3
"""Bygger BRE_Styremote_Finansiering_2026_18.pptx rent fra rev17 + alle endringer."""
import copy, shutil, os
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.opc.package import Part
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI

BASE="media/inbound/openclaw-staged-5a6d151d-e69c-4fb9-9ad1-95fedbd9e7e0/BRE_Styremote_Finansiering_2026_17---69fb27f5-d34a-4722-8712-6dac6ddc4537.pptx"
SENS="media/inbound/openclaw-staged-f1937f2e-3349-417e-b6ef-b620f63b7f06/2026_BRE_Digital_Strategi_Komplett_-_Skrivebeskyttet---c074d2f3-014d-47d4-a1d5-9feea3db0cbe.pptx"
OUT="BRE_Styremote_Finansiering_2026_18.pptx"
shutil.copy(BASE, OUT)
prs=Presentation(OUT)

def variants(s):
    v=[s]
    if " " in s: v.append(s.replace(" "," "))
    return v
def rep_para(para, old, new):
    for ov in variants(old):
        for r in para.runs:
            if ov in r.text: r.text=r.text.replace(ov,new); return True
        full="".join(r.text for r in para.runs)
        if ov in full and para.runs:
            para.runs[0].text=full.replace(ov,new)
            for r in para.runs[1:]: r.text=""
            return True
    return False
def rep_slide(idx, pairs):
    s=prs.slides[idx]; n=0
    for sh in s.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for old,new in pairs:
                    if rep_para(para,old,new): n+=1
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for old,new in pairs:
                            if rep_para(para,old,new): n+=1
    return n

# ---- PROGNOSE + reconciliation (per slide, 0-basert: slide3=2 osv.) ----
rep_slide(2,[("Basisplan 2026: 23,2 MNOK — prognose 24,1","Basisplan 2026: 23,2 MNOK — prognose 26,3 (HubSpot-pipeline)"),
             ("Kilde: PowerOffice resultatregnskap og balanse per 30.06.2026.",
              "Kilde: PowerOffice resultatregnskap og balanse per 30.06.2026. Omsetning hittil 2026 (per 09.08): 16,4 MNOK (live PowerOffice).")])
rep_slide(3,[("24 MNOK","26,3 MNOK"),("prognose 2026 → offensivt løp 2030","prognose 2026 (HubSpot) → offensivt løp 2030"),
             ("3,1x","1,5x"),("årstakt — utviklerens oppdrag","årstakt — utviklerens oppdrag (finansiell plan ~15 %; strategisk ambisjon 35 % SaaS)")])
rep_slide(4,[("2026: prognose 24,1 MNOK — dobling av 2024.","2026: prognose 26,3 MNOK (HubSpot-pipeline) — 2,3× 2024."),
             ("2026P: PowerOffice-prognose.","2026P: prognose (HubSpot-pipeline); budsjett 23,2."),
             ("2027–2030: offensivt løp.","2027–2030: offensivt løp (63 MNOK i 2030); strategisk basismål 55 MNOK.")])
rep_slide(5,[("salgseffekten lagt til engineering).","salgseffekten lagt til engineering). 2026-omsetningsprognose ~26,3 MNOK (HubSpot).")])
rep_slide(7,[("2026 er det stramme året: 10,8 % med selger fra oktober.",
              "2026: prognose ~19 % margin (H1 leverte 14,3 %); selger fra oktober strammer H2.")])
rep_slide(13,[("Marginen ligger på 16–19 % fra 2027 — historisk toppsjikt-nivå.",
               "Prognose 2026 (HubSpot): inntekt ~26,3 MNOK, EBITDA-margin ~19 %, resultat før skatt ~2,8 MNOK — mot budsjettets ±0. Marginen ligger på 16–19 % fra 2027 — historisk toppsjikt-nivå.")])
rep_slide(14,[("Arket «Balanse 2021–2030» i modellen.","Arket «Balanse 2021–2030» i modellen. 2026-kolonnen = budsjett; omsetningsprognose ~26,3 MNOK.")])
# ---- WORST CASE ut + prognose-finansiering ----
rep_slide(16,[("4,4 MNOK","~6,5 MNOK"),("Bank 31.12.2026 (budsjett)","Bank 31.12.2026 (prognose)"),
              ("Samme i stress (null salgseffekt) — basisveksten bærer kassen","Prognosebasert estimat (26,3 MNOK / 19 %) — mer buffer enn planlagt"),
              ("0,3 MNOK","16,4 MNOK"),("Bank i worst case 2029","Omsetning hittil 2026"),
              ("Null salgseffekt + halvert vekst: negativ i 2030 hvis ansettelsene ikke bremses — planen er selvbremsende","Live PowerOffice per 09.08 — ~62 % av prognosen levert"),
              ("3,1x","~1,5x"),("Gjeldsgrad på topp (2026)","Gjeldsgrad topp 2026 (prognose)")])
rep_slide(18,[("Bankbeholdningen 2026–2030: tre baner","Bankbeholdningen 2026–2030: to baner"),
              ("Budsjett- og stressbanen holder seg trygt over vannlinjen — worst case krever ansettelsesbrems.",
               "Forsiktig planbane og stress holder seg trygt over vannlinjen — prognosen ligger over (bank 2026 ~6,5 MNOK).")])
rep_slide(19,[("Worst case-banene er ikke teoretiske.","Stress-scenariet er ikke teoretisk.")])
rep_slide(20,[("3,1x","1,5x")])
# slide19: fjern "Gul = worst case"-avsnitt + chart-serie
s19=prs.slides[18]
for sh in list(s19.shapes):
    if sh.has_text_frame:
        for para in list(sh.text_frame.paragraphs):
            if para.text.strip().startswith("Gul = worst case"):
                para._p.getparent().remove(para._p)
    if sh.has_chart:
        cs=sh.chart._chartSpace
        for ser in list(cs.findall('.//'+qn('c:ser'))):
            tx=ser.find(qn('c:tx')); nm=""
            if tx is not None:
                v=tx.find('.//'+qn('c:v'))
                if v is not None and v.text: nm=v.text
            if "worst case" in nm.lower(): ser.getparent().remove(ser)
print("Tekst-endringer ferdig.")

# ===================== ORG CHART (native) =====================
BLA=RGBColor(0x00,0x56,0x89); MID=RGBColor(0x00,0x92,0xD2); LYS=RGBColor(0x59,0xC2,0xEA)
GUL=RGBColor(0xFA,0xE1,0x00); HVIT=RGBColor(0xFF,0xFF,0xFF); MORK=RGBColor(0x22,0x2B,0x33)
GRA=RGBColor(0xEE,0xF2,0xF5); GRAA=RGBColor(0x88,0x93,0x9b); LB=RGBColor(0xDD,0xEE,0xF6)
def blank_slide():
    lay=min(prs.slide_layouts,key=lambda L:len(L.placeholders))
    s=prs.slides.add_slide(lay)
    for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
    return s
def orgdraw(s):
    def dl(sh): ln=sh.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    def box(l,t,w,h,fill=None,line=None,lw=1.0,dsh=False):
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
        if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
        else: sh.fill.background()
        if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
        else: sh.line.fill.background()
        if dsh: dl(sh)
        sh.shadow.inherit=False; return sh
    def tx(l,t,w,h,text,size=12,color=MORK,bold=False,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE):
        tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
        for i,ln in enumerate(text.split("\n")):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
            r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Calibri"
    def names(x,t,w,h,items):
        box(x,t,w,h,fill=GRA)
        tb=s.shapes.add_textbox(Inches(x+0.12),Inches(t+0.08),Inches(w-0.24),Inches(h-0.16)); tf=tb.text_frame; tf.word_wrap=True
        for j,n in enumerate(items):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
            r=p.add_run(); r.text="• "+n; r.font.size=Pt(11); r.font.color.rgb=MORK; r.font.name="Calibri"; p.space_after=Pt(2)
    box(0,0,13.333,0.9,fill=BLA)
    tx(0.5,0.1,11.5,0.38,"BRE Digital — Organisasjon per august 2026",21,HVIT,True,PP_ALIGN.LEFT)
    tx(0.5,0.52,11.5,0.3,"Etter siste ansettelser (★) + hvor vi styrker fremover (stiplet)",12,LYS,False,PP_ALIGN.LEFT)
    box(5.17,0.98,3.0,0.46,fill=BLA); tx(5.17,0.98,3.0,0.46,"CEO — Frode Lillebakk",12.5,HVIT,True)
    lt=1.52; box(0.55,lt,12.0,0.72,fill=LB); tx(0.68,lt+0.03,1.9,0.66,"LEDERTEAM",10.5,BLA,True,PP_ALIGN.LEFT)
    for i,(role,nm) in enumerate([("CEO","Frode Lillebakk"),("Salg","Roger Øverland ★"),("Leveranse / PL","F. Johansen"),("Utvikling","(lead avklares)")]):
        x=2.3+i*(2.48+0.07); box(x,lt+0.09,2.48,0.54,fill=HVIT,line=MID,lw=1.25)
        tx(x+0.05,lt+0.1,2.38,0.26,role,9,MID,True); tx(x+0.05,lt+0.34,2.38,0.26,nm,10.5,MORK,True)
    st=2.4; tx(0.68,st,2.0,0.5,"STAB (ekstern)",10.5,GRAA,True,PP_ALIGN.LEFT)
    box(3.7,st,2.6,0.5,fill=HVIT,line=GRAA,lw=1.75,dsh=True); tx(3.7,st,2.6,0.5,"Personal — Admento",11.5,MORK,True)
    box(6.5,st,2.6,0.5,fill=HVIT,line=GRAA,lw=1.75,dsh=True); tx(6.5,st,2.6,0.5,"Økonomi — Admento",11.5,MORK,True)
    hy=3.05; hh=0.42
    cols=[("Salg & marked",["Roger Øverland (CSO) ★"]),("Utvikling",["E. Sæther","A.K Sylte","Even Krakeli ★"]),
          ("Drift / leveranse (inkl. support)",["J.T Kallmyr","V. Høgstøyl","Support – kundesaker/tickets"]),
          ("Installasjon",["J.O Brevik","C.J Skotheimsvik","V. Johansen","A.H Bullgård (tavle)"])]
    for i,(head,nm) in enumerate(cols):
        x=0.55+i*(2.85+0.17); box(x,hy,2.85,hh,fill=MID); tx(x+0.05,hy,2.75,hh,head,11,HVIT,True); names(x,hy+0.47,2.85,1.5,nm)
    yb=5.12; box(0.55,yb-0.1,12.0,0.03,fill=GUL); tx(0.55,yb,12.0,0.32,"HVOR VI STYRKER OSS FREMOVER",12.5,BLA,True,PP_ALIGN.LEFT)
    for i,(h,note) in enumerate([("Økonomi / Controller","fra Admento → egen controller"),("Software #2","2028 — cloud-dybde"),("Innesalg / tilbud","mer solgt per krone"),("Fagkapasitet","installasjon, 2028→")]):
        x=0.55+i*(2.87+0.19); ys=yb+0.38; box(x,ys,2.87,0.98,fill=HVIT,line=GUL,lw=2.25,dsh=True)
        tx(x+0.06,ys+0.09,2.75,0.45,"＋ "+h,12,BLA,True); tx(x+0.06,ys+0.52,2.75,0.4,note,9.5,GRAA,False)
    tx(0.55,6.55,12.0,0.3,"★ = nyansatt. F. Johansen frigjøres til prosjektleder for leveranse. Personal og økonomi ivaretas eksternt av Admento. Stiplet = eksternt/planlagt.",9,GRAA,False,PP_ALIGN.LEFT)

org=blank_slide(); orgdraw(org)
# ===================== HANDLINGSPLAN (kopi fra Sensacon idx 28) =====================
sens=Presentation(SENS); src_slide=sens.slides[28]
hp=blank_slide()
pkg=prs.part.package; existing=set(str(p.partname) for p in pkg.iter_parts()); idmap={}
def newrid(old):
    if old in idmap: return idmap[old]
    img=src_slide.part.related_part(old); ext=img.partname.ext or "png"; n=1
    while f"/ppt/media/hpimg_{n}.{ext}" in existing: n+=1
    pn=f"/ppt/media/hpimg_{n}.{ext}"; existing.add(pn)
    part=Part(PackURI(pn),img.content_type,pkg,img.blob); rid=hp.part.relate_to(part,RT.IMAGE); idmap[old]=rid; return rid
for shp in src_slide.shapes: hp.shapes._spTree.append(copy.deepcopy(shp._element))
for el in hp.shapes._spTree.iter():
    o=el.get(qn('r:embed'))
    if o: el.set(qn('r:embed'),newrid(o))

# ---- rekkefølge: org etter slide13 (idx13), handlingsplan etter org (idx14) ----
lst=prs.slides._sldIdLst; ids=list(lst)
hp_node=ids[-1]; org_node=ids[-2]
lst.remove(org_node); lst.remove(hp_node)
lst.insert(13, org_node); lst.insert(14, hp_node)
prs.save(OUT)
print("Ferdig. Slides:",len(prs.slides._sldIdLst))
EOF
